from io import BytesIO
from typing import BinaryIO

import fitz  # PyMuPDF
from docx import Document
from pptx import Presentation


def read_pdf(file: BinaryIO) -> str:
    """Read a PDF file-like object using PyMuPDF and return extracted text."""
    # PyMuPDF expects a bytes object or a file path, so we read the bytes
    file_bytes = file.read()
    doc = fitz.open(stream=file_bytes, filetype="pdf")
    text_parts: list[str] = []
    for page in doc:
        text_parts.append(page.get_text())
    doc.close()
    return "\n".join(text_parts).strip()


def read_docx(file: BinaryIO) -> str:
    """Read a DOCX file-like object using python-docx; extracts paragraphs and tables."""
    file.seek(0)
    doc = Document(file)
    parts: list[str] = []
    for p in doc.paragraphs:
        if p.text.strip():
            parts.append(p.text)
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                parts.append(row_text)
    return "\n".join(parts).strip()


def read_pptx(file: BinaryIO) -> str:
    """Read a PPTX file-like object; extracts text from slides and tables."""
    file.seek(0)
    file_bytes = file.read()
    prs = Presentation(BytesIO(file_bytes))
    parts: list[str] = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    row_text = " | ".join(
                        cell.text_frame.text.strip()
                        for cell in row.cells
                        if cell.text_frame.text.strip()
                    )
                    if row_text:
                        parts.append(row_text)
            elif shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    if para.text.strip():
                        parts.append(para.text)
    return "\n".join(parts).strip()


def read_txt(file: BinaryIO, encoding: str = "utf-8") -> str:
    """Read a plain text file-like object and return text."""
    # Предполагаем, что это бинарный файловый объект (как UploadFile.file)
    file.seek(0)
    content = file.read()
    if isinstance(content, bytes):
        return content.decode(encoding, errors="ignore").strip()
    return str(content).strip()


def get_word_count(text: str) -> int:
    """Return a simple word count for the given text."""
    if not text:
        return 0
    # Простое разбиение по пробелам/переводам строк
    words = text.split()
    return len(words)
